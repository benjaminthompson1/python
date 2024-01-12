# Required Libraries:
# python-pptx for handling PowerPoint files: Install via `pip install python-pptx`
# python-dotenv for loading environment variables: Install via `pip install python-dotenv`
# GenAI specific libraries for AI model interaction: Installation instructions should be provided by GenAI

import time, os
from dotenv import load_dotenv
from genai.credentials import Credentials
from genai.model import Model
from genai.schemas import GenerateParams
from pptx import Presentation

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return " ".join(text)

load_dotenv()
api_key = os.getenv("GENAI_KEY", None)
api_url = os.getenv("GENAI_API", None)
creds = Credentials(api_key, api_endpoint=api_url)

print("\n------------- Interactive AI Chat -------------\n")

# Initialize AI model (Alice)
alice_params = GenerateParams(decoding_method="sample", max_new_tokens=500, temperature=0.15)
alice = Model("google/flan-t5-xxl", params=alice_params, credentials=creds)

# Prompt for PowerPoint file path
pptx_file_path = input("Enter the path to your PowerPoint file: ")
pptx_text = extract_text_from_pptx(pptx_file_path)

# Debug option
debug_mode = input("Enable debug mode? (yes/no): ").lower() == "yes"
if debug_mode:
    print("\n[Debug] Extracted PowerPoint Text: ")
    print(pptx_text)
    print("\n")

while True:
    user_input = input("Your message (type 'QUIT' to exit): ")
    if user_input == "QUIT":
        break

    # Concatenate user input with the PowerPoint text for context
    combined_input = pptx_text + " " + user_input

    alice_response = alice.generate([combined_input])
    alice_gen = alice_response[0].generated_text
    print(f"[Alice] --> {alice_gen}\n")

    time.sleep(0.5)
