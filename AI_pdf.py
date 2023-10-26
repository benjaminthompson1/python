import os
from dotenv import load_dotenv
from genai.extensions.langchain import LangChainInterface
from genai.schemas import GenerateParams
from genai.credentials import Credentials
from langchain.document_loaders import PyPDFLoader
from pptx import Presentation

MAX_TOKENS = 4096

def setup_genai_credentials():
    load_dotenv()
    api_key = os.getenv("GENAI_KEY")
    api_url = os.getenv("GENAI_API")
    
    if not api_key or not api_url:
        print("Error: Ensure the GENAI_KEY and GENAI_API are set in the environment or .env file.")
        exit()

    return Credentials(api_key, api_endpoint=api_url)

def extract_text_from_pdf(pdf_path):
    loader = PyPDFLoader(pdf_path)
    pdf_content = loader.load()
    try:
        return ' '.join([doc.page_content for doc in pdf_content])
    except AttributeError:
        print("Error: Failed to extract text from the Document objects.")
        available_attributes = dir(pdf_content[0]) if pdf_content else []
        print(f"Available attributes/methods: {available_attributes}")
        exit()

def extract_text_from_powerpoint(ppt_path):
    prs = Presentation(ppt_path)
    return ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]).strip()

def extract_text_from_document(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()
    extractors = {
        '.pdf': extract_text_from_pdf,
        '.ppt': extract_text_from_powerpoint,
        '.pptx': extract_text_from_powerpoint
    }

    return extractors.get(file_extension, unsupported_format)(file_path)

def unsupported_format(file_path):
    print(f"Unsupported file format: {os.path.splitext(file_path)[1].lower()}")
    exit()

def interact_with_user(extracted_text, langchain_model):
    print("\n---------- Ask questions about the document ----------\n")

    while True:
        question = input("Ask a question about the document (or type 'exit' to quit): ")
        if question.lower() == 'exit':
            break
        
        tokens = extracted_text.split()  # Simple tokenization
        truncated_content = ' '.join(tokens[-MAX_TOKENS:])
        context_question = f"Based on this content from the document: {truncated_content}\n{question}"

        response = langchain_model(context_question)
        print(f"Answer: {response}\n")

def main():
    print("\n------------- Example (LangChain) -------------\n")
    creds = setup_genai_credentials()

    file_path = input("Enter the path to the document file (PDF or PowerPoint): ")
    extracted_text = extract_text_from_document(file_path)

    # Ask the user if they want to view the document content
    while True:
        view_content = input("Would you like to view the contents of the document? (yes/no): ").lower()
        if view_content in ["yes", "no"]:
            break
        print("Please enter 'yes' or 'no'.")

    if view_content == "yes":
        print("\n------ Begin Document Content ------\n")
        print(extracted_text)
        print("\n------ End Document Content ------\n")

    print("Using GenAI Model expressed as LangChain Model via LangChainInterface:")
    params = GenerateParams(decoding_method="greedy")
    langchain_model = LangChainInterface(model="google/flan-t5-xxl", params=params, credentials=creds)

    interact_with_user(extracted_text, langchain_model)
    print("Exiting...")

if __name__ == "__main__":
    main()
